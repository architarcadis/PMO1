"""
Export Module for PMO Pulse Dashboard

Provides functions to export data and visualizations to Excel and PowerPoint.
"""

import pandas as pd
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

def export_to_excel(projects_df, tasks_df, risks_df, history_df):
    """
    Export all project data to an Excel file.
    
    Parameters:
    -----------
    projects_df : pandas.DataFrame
        DataFrame containing project information
    tasks_df : pandas.DataFrame
        DataFrame containing task information
    risks_df : pandas.DataFrame
        DataFrame containing risk information
    history_df : pandas.DataFrame
        DataFrame containing historical performance data
        
    Returns:
    --------
    bytes
        Excel file as bytes for download
    """
    output = io.BytesIO()
    
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        # Export projects
        projects_df.to_excel(writer, sheet_name='Projects', index=False)
        
        # Export tasks
        tasks_df.to_excel(writer, sheet_name='Tasks', index=False)
        
        # Export risks
        risks_df.to_excel(writer, sheet_name='Risks', index=False)
        
        # Export performance history
        history_df.to_excel(writer, sheet_name='Performance History', index=False)
    
    output.seek(0)
    return output.getvalue()

def generate_data_template():
    """
    Generate an Excel template with the required columns and format for data upload
    
    Returns:
    --------
    bytes
        Excel template file as bytes for download
    """
    output = io.BytesIO()
    
    # Create empty dataframes with the required column structure
    projects_template = pd.DataFrame(columns=[
        'project_id', 'project_name', 'description', 'sector', 'project_manager',
        'planned_start', 'planned_end', 'actual_start', 'actual_end', 
        'budget', 'actual_cost', 'status', 'spi', 'cpi',
        'budget_variance_pct', 'schedule_variance_days', 'completion_pct',
        'priority', 'location', 'latitude', 'longitude'
    ])
    
    tasks_template = pd.DataFrame(columns=[
        'task_id', 'project_id', 'task_name', 'description', 
        'planned_start', 'planned_end', 'actual_start', 'actual_end',
        'planned_cost', 'actual_cost', 'earned_value', 'resource_id',
        'is_milestone', 'is_critical', 'predecessor_ids', 'status',
        'completion_pct', 'total_float'
    ])
    
    risks_template = pd.DataFrame(columns=[
        'risk_id', 'project_id', 'title', 'description', 'impact',
        'probability', 'severity', 'category', 'owner', 'mitigation_plan',
        'contingency_plan', 'status', 'identified_date', 'last_update',
        'impact_areas'
    ])
    
    resources_template = pd.DataFrame(columns=[
        'resource_id', 'resource_name', 'role', 'department', 'cost_rate',
        'availability', 'skills', 'allocation_percentage', 'email',
        'assigned_projects'
    ])
    
    # Add sample data and format guidelines to help users
    sample_project = {
        'project_id': 'PROJ001',
        'project_name': 'Example Project',
        'description': 'This is a sample project entry',
        'sector': 'Infrastructure',
        'project_manager': 'Jane Smith',
        'planned_start': '2025-01-01',
        'planned_end': '2025-06-30',
        'actual_start': '2025-01-05',
        'actual_end': None,
        'budget': 500000,
        'actual_cost': 210000,
        'status': 'In Progress',
        'spi': 0.95,
        'cpi': 1.05,
        'budget_variance_pct': 5,
        'schedule_variance_days': -3,
        'completion_pct': 40,
        'priority': 'High',
        'location': 'New York',
        'latitude': 40.7128,
        'longitude': -74.0060
    }
    
    sample_task = {
        'task_id': 'TASK001',
        'project_id': 'PROJ001',
        'task_name': 'Example Task',
        'description': 'This is a sample task entry',
        'planned_start': '2025-01-15',
        'planned_end': '2025-02-15',
        'actual_start': '2025-01-16',
        'actual_end': None,
        'planned_cost': 50000,
        'actual_cost': 25000,
        'earned_value': 20000,
        'resource_id': 'RES001',
        'is_milestone': False,
        'is_critical': True,
        'predecessor_ids': '',
        'status': 'In Progress',
        'completion_pct': 50,
        'total_float': 0
    }
    
    sample_risk = {
        'risk_id': 'RISK001',
        'project_id': 'PROJ001',
        'title': 'Example Risk',
        'description': 'This is a sample risk entry',
        'impact': 'High',
        'probability': 'Medium',
        'severity': 'High',
        'category': 'Schedule',
        'owner': 'John Doe',
        'mitigation_plan': 'Add backup resources',
        'contingency_plan': 'Extend deadline by 2 weeks',
        'status': 'Active',
        'identified_date': '2025-01-05',
        'last_update': '2025-01-20',
        'impact_areas': 'Schedule, Budget'
    }
    
    sample_resource = {
        'resource_id': 'RES001',
        'resource_name': 'John Smith',
        'role': 'Engineer',
        'department': 'Engineering',
        'cost_rate': 75,
        'availability': 80,
        'skills': 'Civil Engineering, CAD, Project Management',
        'allocation_percentage': 60,
        'email': 'john.smith@example.com',
        'assigned_projects': 'PROJ001, PROJ002'
    }
    
    # Add sample rows to the templates
    projects_template.loc[0] = sample_project
    tasks_template.loc[0] = sample_task
    risks_template.loc[0] = sample_risk
    resources_template.loc[0] = sample_resource
    
    # Create the Excel file with multiple sheets
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        projects_template.to_excel(writer, sheet_name='Projects', index=False)
        tasks_template.to_excel(writer, sheet_name='Tasks', index=False)
        risks_template.to_excel(writer, sheet_name='Risks', index=False)
        resources_template.to_excel(writer, sheet_name='Resources', index=False)
        
        # Add an instructions sheet
        instructions_df = pd.DataFrame({
            'Field': [
                'Date formats',
                'IDs',
                'Status values',
                'Priority values',
                'Completion percentage',
                'Required fields',
                'Relationship between sheets'
            ],
            'Description': [
                'Use YYYY-MM-DD format for all dates',
                'IDs should be unique and consistent across sheets',
                'Use consistent status values like: "Not Started", "In Progress", "Completed", "On Hold", "At Risk"',
                'Use values: "Low", "Medium", "High", "Critical"',
                'Enter as whole numbers (0-100)',
                'project_id, project_name, planned_start, planned_end, and budget are required for projects',
                'Foreign keys must match: project_id in Tasks and Risks must exist in Projects sheet'
            ]
        })
        instructions_df.to_excel(writer, sheet_name='Instructions', index=False)
    
    return output.getvalue()

def export_to_powerpoint(projects_df, tasks_df, risks_df, history_df, portfolio_kpis):
    """
    Export portfolio overview to a PowerPoint presentation.
    
    Parameters:
    -----------
    projects_df : pandas.DataFrame
        DataFrame containing project information
    tasks_df : pandas.DataFrame
        DataFrame containing task information
    risks_df : pandas.DataFrame
        DataFrame containing risk information
    history_df : pandas.DataFrame
        DataFrame containing historical performance data
    portfolio_kpis : dict
        Dictionary containing portfolio KPIs
        
    Returns:
    --------
    bytes
        PowerPoint file as bytes for download
    """
    # Create presentation
    prs = Presentation()
    
    # Define Arcadis colors
    ARCADIS_ORANGE = RGBColor(230, 115, 0)
    ARCADIS_BLACK = RGBColor(0, 0, 0)
    ARCADIS_GREY = RGBColor(108, 117, 125)
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "PMO Pulse Dashboard Report"
    subtitle.text = f"Generated on {datetime.datetime.now().strftime('%B %d, %Y')}"
    
    # Adjust title formatting
    title.text_frame.paragraphs[0].font.color.rgb = ARCADIS_ORANGE
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Add portfolio overview slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Portfolio Overview"
    
    # Add portfolio KPIs
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = f"Total Projects: {len(projects_df)}"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = f"Portfolio SPI: {portfolio_kpis['avg_spi']:.2f}"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = f"Portfolio CPI: {portfolio_kpis['avg_cpi']:.2f}"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = f"Active Projects: {portfolio_kpis['active_projects']}"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = f"Projects At Risk: {portfolio_kpis['at_risk_projects']}"
    p.font.size = Pt(18)
    
    # Add projects by status slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Projects by Status"
    
    # Add status breakdown
    content = slide.placeholders[1]
    tf = content.text_frame
    
    status_counts = projects_df['status'].value_counts()
    
    for status, count in status_counts.items():
        p = tf.add_paragraph()
        p.text = f"{status}: {count} projects"
        p.font.size = Pt(16)
    
    # Add projects by sector slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Projects by Sector"
    
    # Add sector breakdown
    content = slide.placeholders[1]
    tf = content.text_frame
    
    sector_counts = projects_df['sector'].value_counts()
    
    for sector, count in sector_counts.items():
        p = tf.add_paragraph()
        p.text = f"{sector}: {count} projects"
        p.font.size = Pt(16)
    
    # Add projects needing attention slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Projects Needing Attention"
    
    # Add list of at-risk projects
    content = slide.placeholders[1]
    tf = content.text_frame
    
    at_risk = projects_df[projects_df['status'].isin(['At Risk', 'Delayed'])]
    
    if not at_risk.empty:
        for _, row in at_risk.iterrows():
            p = tf.add_paragraph()
            p.text = f"{row['project_name']} - {row['status']}"
            p.font.size = Pt(14)
            p.level = 0
            
            # Add bullet point with manager
            p = tf.add_paragraph()
            p.text = f"Manager: {row['project_manager']}"
            p.font.size = Pt(12)
            p.level = 1
    else:
        p = tf.add_paragraph()
        p.text = "No projects currently at risk."
        p.font.size = Pt(16)
    
    # Save presentation to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output.getvalue()
